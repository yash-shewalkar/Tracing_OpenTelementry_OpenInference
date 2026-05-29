import os
import logging
import argparse
import tempfile
from typing import Optional

from openai import OpenAI
from markitdown import MarkItDown
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S"
)
logger = logging.getLogger(__name__)

def create_groq_client() -> Optional[OpenAI]:
    """Initializes the OpenAI client with Groq's base URL and API key."""
    api_key = os.getenv("GROQ_API_KEY")
    if not api_key:
        logger.warning("GROQ_API_KEY not found. Please set the environment variable.")
        return None
    try:
        client = OpenAI(
            api_key=api_key,
            base_url="https://api.groq.com/openai/v1",
        )
        logger.info("Groq client initialized successfully.")
        return client
    except Exception as e:
        logger.exception(f"Failed to initialize Groq client: {e}")
        return None

def clean_shapes_recursive(shapes, slide_num: int) -> bool:
    """
    Recursively scans shapes (and grouped shapes) to find and remove 
    broken images that throw 'no embedded image' ValueErrors.
    """
    needs_saving = False
    shapes_to_remove = []
    
    for shape in shapes:
        # 1. If it's a group, recurse into it
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            if clean_shapes_recursive(shape.shapes, slide_num):
                needs_saving = True
                
        # 2. Safely check if the shape class supports the 'image' property.
        # We check type(shape) so we don't accidentally execute the property descriptor.
        if hasattr(type(shape), "image"):
            try:
                # Deliberately access the property to see if it crashes
                _ = shape.image
            except ValueError as e:
                # Catch the specific python-pptx crash
                if "no embedded image" in str(e).lower():
                    shapes_to_remove.append(shape)
                    needs_saving = True
            except Exception:
                # Ignore other exceptions (like UnidentifiedImageError) to let MarkItDown handle them
                pass
                
    # 3. Remove the broken shapes from the XML tree
    for shape in shapes_to_remove:
        logger.warning(f"Removing broken image shape on Slide {slide_num}")
        element = shape._element
        element.getparent().remove(element)
        
    return needs_saving

def sanitize_pptx(pptx_path: str) -> str:
    """
    Scans the PPTX for broken or empty image shapes that crash MarkItDown
    and removes them. Returns the path to the clean file.
    """
    logger.info("Scanning PPTX for broken/linked images...")
    try:
        prs = Presentation(pptx_path)
        needs_saving = False
        
        for slide_num, slide in enumerate(prs.slides, 1):
            if clean_shapes_recursive(slide.shapes, slide_num):
                needs_saving = True
                
        if needs_saving:
            # Save the cleaned presentation to a temporary file
            temp_fd, temp_path = tempfile.mkstemp(suffix=".pptx")
            os.close(temp_fd) # Close file descriptor, python-pptx will handle writing
            prs.save(temp_path)
            logger.info(f"Created sanitized temporary PPTX file: {temp_path}")
            return temp_path
            
        logger.info("PPTX is clean. No broken images found.")
        return pptx_path
        
    except Exception as e:
        logger.error(f"Failed to sanitize PPTX. Proceeding with original. Error: {e}")
        return pptx_path

def convert_pptx_to_md(pptx_path: str, output_md_path: str, vision_model: str = "meta-llama/llama-4-scout-17b-16e-instruct") -> bool:
    """
    Converts a PPTX file to Markdown, using Groq LLM to describe embedded images.
    """
    if not os.path.exists(pptx_path):
        logger.error(f"Input file does not exist: {pptx_path}")
        return False

    # 1. Preprocess the file to prevent "no embedded image" crashes
    safe_pptx_path = sanitize_pptx(pptx_path)
    
    # 2. Setup the LLM client
    client = create_groq_client()
    if client:
        logger.info(f"Using LLM for image descriptions. Model: {vision_model}")
        md = MarkItDown(llm_client=client, llm_model=vision_model)
    else:
        logger.warning("Falling back to text-only conversion without LLM image descriptions.")
        md = MarkItDown()

    # 3. Convert via MarkItDown
    try:
        logger.info(f"Starting conversion for: {pptx_path}")
        result = md.convert(safe_pptx_path)
        
        with open(output_md_path, "w", encoding="utf-8") as f:
            f.write(result.text_content)
            
        logger.info(f"Conversion successful! Markdown saved to: {output_md_path}")
        success = True

    except Exception as e:
        logger.exception(f"An error occurred during PPTX conversion: {e}")
        success = False
        
    finally:
        # Cleanup the temporary sanitized file if one was created
        if safe_pptx_path != pptx_path and os.path.exists(safe_pptx_path):
            os.remove(safe_pptx_path)
            logger.info("Cleaned up temporary sanitized file.")

    return success

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Convert PPTX to Markdown with LLM image descriptions via Groq.")
    parser.add_argument("input_file", help="Path to the input .pptx file")
    parser.add_argument("-o", "--output", help="Path to the output .md file (optional)", default=None)
    parser.add_argument("--model", help="Groq Vision model to use", default="meta-llama/llama-4-scout-17b-16e-instruct")
    
    args = parser.parse_args()
    
    out_path = args.output
    if not out_path:
        base_name = os.path.splitext(args.input_file)[0]
        out_path = f"{base_name}.md"
        
    convert_pptx_to_md(args.input_file, out_path, vision_model=args.model)